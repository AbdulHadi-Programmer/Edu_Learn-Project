# import streamlit as st
# import PyPDF2
# import docx
# import nltk
# from nltk.tokenize import sent_tokenize
# from sumy.parsers.plaintext import PlaintextParser
# from sumy.nlp.tokenizers import Tokenizer
# from sumy.summarizers.lsa import LsaSummarizer
# from sumy.summarizers.lex_rank import LexRankSummarizer
# import pyperclip
# from pptx import Presentation
# import random

# # Download required NLTK resources
# nltk.download('punkt')

# # Read PPTX
# def read_pptx(file):
#     prs = Presentation(file)
#     text = "\n".join([
#         shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
#     ])
#     return text

# # Read PDF
# def read_pdf(file):
#     pdf_reader = PyPDF2.PdfReader(file)
#     text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
#     return text

# # Read DOCX
# def read_docx(file):
#     doc = docx.Document(file)
#     text = "\n".join([para.text for para in doc.paragraphs])
#     return text

# # Read TXT
# def read_txt(file):
#     return file.read().decode("utf-8")

# # Extractive Summary
# def extractive_summary(text, num_sentences=5):
#     sentences = sent_tokenize(text)
#     return " ".join(sentences[:num_sentences])

# # Abstractive Summary (LSA / LexRank)
# def abstractive_summary(text, num_sentences=5, method="lsa"):
#     parser = PlaintextParser.from_string(text, Tokenizer("english"))
#     summarizer = LsaSummarizer() if method == "lsa" else LexRankSummarizer()
#     summary = summarizer(parser.document, num_sentences)
#     return " ".join(str(sentence) for sentence in summary)

# # MCQ Generator (as-is from your version)
# def generate_mcqs(text, num_questions=3):
#     sentences = sent_tokenize(text)
#     if len(sentences) < num_questions:
#         num_questions = len(sentences)
#     selected_sentences = random.sample(sentences, num_questions)
#     questions = []
#     for sentence in selected_sentences:
#         words = sentence.split()
#         if len(words) < 4:
#             continue
#         blank_index = random.randint(0, len(words) - 1)
#         answer = words[blank_index]
#         words[blank_index] = "_____"
#         question = " ".join(words)
#         options = [answer] + random.sample([w for w in words if w != "_____" and w.isalpha()], min(3, len(words)-1))
#         random.shuffle(options)
#         questions.append({
#             "question": question,
#             "options": options,
#             "answer": answer
#         })
#     return questions

# # Streamlit App
# def main():
#     st.title("🧠 AI Summarizer + MCQ Generator")

#     uploaded_file = st.file_uploader("📂 Upload PDF, DOCX, TXT, or PPTX", type=["pdf", "docx", "txt", "pptx"])

#     if uploaded_file:
#         file_type = uploaded_file.type
#         if "pdf" in file_type:
#             text = read_pdf(uploaded_file)
#         elif "word" in file_type or "docx" in uploaded_file.name:
#             text = read_docx(uploaded_file)
#         elif "text" in file_type:
#             text = read_txt(uploaded_file)
#         elif "pptx" in uploaded_file.name:
#             text = read_pptx(uploaded_file)
#         else:
#             st.error("❌ Unsupported file format")
#             return

#         # --- Summary Section ---
#         st.header("📚 Summarization")

#         summary_type = st.radio("🔍 Choose summary type:", ["Extractive", "Abstractive"])
#         num_sentences = st.slider("✂ Select number of sentences:", 3, 20, 5)

#         if summary_type == "Extractive":
#             summary = extractive_summary(text, num_sentences)
#         else:
#             method = st.radio("📌 Abstractive Method:", ['lsa', 'lexrank'])
#             st.info("LSA: Best for structured academic text\nLexRank: Better for long or unstructured content")
#             summary = abstractive_summary(text, num_sentences, method)

#         st.subheader("📝 Summary")
#         st.text_area("Summary Output:", summary, height=200)

#         if st.button("📋 Copy Summary"):
#             pyperclip.copy(summary)
#             st.success("✅ Copied to clipboard!")

#         if st.button("💾 Download Summary"):
#             with open("summary.txt", "w", encoding="utf-8") as f:
#                 f.write(summary)
#             st.download_button(label="📥 Download", data=summary, file_name="summary.txt", mime="text/plain")

#         # --- MCQ Section ---
#         st.header("❓ MCQ Generator")

#         if st.button("🎯 Generate MCQs"):
#             mcqs = generate_mcqs(text, _numquestions=5)
#             if mcqs:
#                 for i, q in enumerate(mcqs, 1):
#                     st.markdown(f"**Q{i}. {q['question']}**")
#                     for opt in q["options"]:
#                         st.markdown(f"- {opt}")
#                     st.markdown(f"**✅ Answer:** {q['answer']} \n---")
#             else:
#                 st.warning("Not enough content to generate MCQs.")

# if __name__ == "__main__":
#     main()

## Perfect but the problem is 
# import streamlit as st
# import PyPDF2
# import docx
# import nltk
# from nltk.tokenize import sent_tokenize
# from sumy.parsers.plaintext import PlaintextParser
# from sumy.nlp.tokenizers import Tokenizer
# from sumy.summarizers.lsa import LsaSummarizer
# from sumy.summarizers.lex_rank import LexRankSummarizer
# import pyperclip
# from pptx import Presentation
# import random

# nltk.download('punkt')

# # Read PPTX
# def read_pptx(file):
#     prs = Presentation(file)
#     text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
#     return text

# # Read PDF
# def read_pdf(file):
#     pdf_reader = PyPDF2.PdfReader(file)
#     text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
#     return text

# # Read DOCX
# def read_docx(file):
#     doc = docx.Document(file)
#     text = "\n".join([para.text for para in doc.paragraphs])
#     return text

# # Read TXT
# def read_txt(file):
#     return file.read().decode("utf-8")

# # Extractive Summary
# def extractive_summary(text, num_sentences=5):
#     sentences = sent_tokenize(text)
#     return " ".join(sentences[:num_sentences])

# # Abstractive Summary (LSA / LexRank)
# def abstractive_summary(text, num_sentences=5, method="lsa"):
#     parser = PlaintextParser.from_string(text, Tokenizer("english"))
#     summarizer = LsaSummarizer() if method == "lsa" else LexRankSummarizer()
#     summary = summarizer(parser.document, num_sentences)
#     return " ".join(str(sentence) for sentence in summary)

# # MCQ Generator
# def generate_mcqs(text, num_questions=3):
#     sentences = sent_tokenize(text)
#     if len(sentences) < num_questions:
#         num_questions = len(sentences)
#     selected_sentences = random.sample(sentences, num_questions)
#     questions = []
#     for sentence in selected_sentences:
#         words = sentence.split()
#         if len(words) < 4:
#             continue
#         blank_index = random.randint(0, len(words) - 1)
#         answer = words[blank_index]
#         words[blank_index] = "_____"
#         question = " ".join(words)
#         options = [answer] + random.sample([w for w in words if w != "_____" and w.isalpha()], min(3, len(words)-1))
#         random.shuffle(options)
#         questions.append({
#             "question": question,
#             "options": options,
#             "answer": answer
#         })
#     return questions

# # Streamlit App
# def main():
#     st.title("🧠 AI Summarizer + MCQ Generator")

#     uploaded_file = st.file_uploader("📂 Upload PDF, DOCX, TXT, or PPTX", type=["pdf", "docx", "txt", "pptx"])

#     if uploaded_file:
#         file_type = uploaded_file.type
#         if "pdf" in file_type:
#             text = read_pdf(uploaded_file)
#         elif "word" in file_type or "docx" in uploaded_file.name:
#             text = read_docx(uploaded_file)
#         elif "text" in file_type:
#             text = read_txt(uploaded_file)
#         elif "pptx" in uploaded_file.name:
#             text = read_pptx(uploaded_file)
#         else:
#             st.error("❌ Unsupported file format")
#             return

#         # --- Summary Section ---
#         st.header("📚 Summarization")

#         summary_type = st.radio("🔍 Choose summary type:", ["Extractive", "Abstractive"])
#         num_sentences = st.slider("✂ Select number of sentences:", 3, 20, 5)

#         if summary_type == "Extractive":
#             summary = extractive_summary(text, num_sentences)
#         else:
#             method = st.radio("📌 Abstractive Method:", ['lsa', 'lexrank'])
#             st.info("LSA: Best for structured academic text\nLexRank: Better for long or unstructured content")
#             summary = abstractive_summary(text, num_sentences, method)

#         st.subheader("📝 Summary")
#         st.text_area("Summary Output:", summary, height=200)

#         if st.button("📋 Copy Summary"):
#             pyperclip.copy(summary)
#             st.success("✅ Copied to clipboard!")

#         if st.button("💾 Download Summary"):
#             with open("summary.txt", "w", encoding="utf-8") as f:
#                 f.write(summary)
#             st.download_button(label="📥 Download", data=summary, file_name="summary.txt", mime="text/plain")

#         # --- MCQ Section ---
#         st.header("❓ MCQ Generator")

#         num_mcqs = st.slider("🔢 Select number of MCQs:", 5, 25, 5)

#         if st.button("🎯 Generate MCQs"):
#             mcqs = generate_mcqs(text, num_questions=num_mcqs)
#             st.session_state["mcqs"] = mcqs
#             st.session_state["show_answers"] = False

#         if "mcqs" in st.session_state:
#             for i, q in enumerate(st.session_state["mcqs"], 1):
#                 st.markdown(f"**Q{i}. {q['question']}**")
#                 st.radio("Choose an option:", q["options"], key=f"mcq_{i}")

#             if st.button("👁️ Show Answers"):
#                 st.session_state["show_answers"] = True

#             if st.session_state.get("show_answers", False):
#                 st.markdown("### ✅ Correct Answers")
#                 for i, q in enumerate(st.session_state["mcqs"], 1):
#                     st.markdown(f"**Q{i} Answer:** {q['answer']}")

# if __name__ == "__main__":
#     main()

## Last Version : (Some changes must to made )
# import streamlit as st
# import PyPDF2
# import docx
# import nltk
# from nltk.tokenize import sent_tokenize
# from sumy.parsers.plaintext import PlaintextParser
# from sumy.nlp.tokenizers import Tokenizer
# from sumy.summarizers.lsa import LsaSummarizer
# from sumy.summarizers.lex_rank import LexRankSummarizer
# import pyperclip
# from pptx import Presentation
# import random

# nltk.download('punkt', quiet=True)

# # --------------------------- File Reading --------------------------- #
# def read_pptx(file):
#     prs = Presentation(file)
#     text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
#     return text

# def read_pdf(file):
#     pdf_reader = PyPDF2.PdfReader(file)
#     text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
#     return text

# def read_docx(file):
#     doc = docx.Document(file)
#     text = "\n".join([para.text for para in doc.paragraphs])
#     return text

# def read_txt(file):
#     return file.read().decode("utf-8")

# # --------------------------- Summarization --------------------------- #
# # def extractive_summary(text, num_sentences):
# #     sentences = sent_tokenize(text)
# #     selected = sentences[:num_sentences]
# #     formatted = []

# #     for sent in selected:
# #         if "*" in sent:
# #             bullets = sent.split("*")
# #             for b in bullets:
# #                 line = b.strip()
# #                 if line:
# #                     formatted.append(f"• {line}")
# #         else:
# #             formatted.append(sent.strip())

# #     return "\n".join(formatted)
# import re

# def preprocess_text(text):
#     # Join broken bullet points like:
#     # 1.\nBudgeting -> 1. Budgeting
#     text = re.sub(r'(\d+)\.\s*\n', r'\1. ', text)
    
#     # Fix lines like "Introduction\nSome sentence..." → "Introduction: Some sentence..."
#     text = re.sub(r'([A-Za-z ]+)\n([A-Z])', r'\1: \2', text)

#     # Remove excessive newlines
#     text = re.sub(r'\n+', ' ', text)

#     return text.strip()

# def extractive_summary(text, num_sentences):
#     # Clean up and fix broken structure
#     clean_text = preprocess_text(text)
    
#     # Tokenize clean version
#     sentences = sent_tokenize(clean_text)
#     selected = sentences[:num_sentences]
    
#     formatted = []
#     for sent in selected:
#         # Bullet everything
#         formatted.append(f"• {sent.strip()}")
    
#     return "\n".join(formatted)

# ############################################################
# def clean_abstractive_output(summary, method):
#     if method == "lsa":
#         # Fix bullet points: put each bullet on a new line
#         summary = summary.replace("*", "\n•")
#         summary = re.sub(r'\n+', '\n', summary)  # collapse extra newlines
#     elif method == "lexrank":
#         # Fix LexRank formatting issues
#         summary = summary.replace('\n', ' ')  # remove hard newlines
#         summary = re.sub(r'\s{2,}', ' ', summary)  # remove extra spaces
#         summary = summary.strip()
#     return summary

# def abstractive_summary(text, num_sentences=5, method="lsa"):
#     parser = PlaintextParser.from_string(text, Tokenizer("english"))
#     summarizer = LsaSummarizer() if method == "lsa" else LexRankSummarizer()
#     summary = summarizer(parser.document, num_sentences)

#     raw_summary = " ".join([str(sentence) for sentence in summary])
#     clean_summary = clean_abstractive_output(raw_summary, method)
#     return clean_summary


# # def abstractive_summary(text, num_sentences=5, method="lsa"):
# #     parser = PlaintextParser.from_string(text, Tokenizer("english"))
# #     summarizer = LsaSummarizer() if method == "lsa" else LexRankSummarizer()
# #     summary = summarizer(parser.document, num_sentences)

# #     # Create a formatted summary like the extractive summary
# #     formatted_summary = " ".join([str(sentence) for sentence in summary])
# #     return formatted_summary

# # --------------------------- MCQ Generator --------------------------- #
# def generate_mcqs(text, num_questions=3):
#     sentences = sent_tokenize(text)
#     random.shuffle(sentences)
#     valid_questions = []

#     for sentence in sentences:
#         words = sentence.split()
#         valid_words = [w for w in words if w.isalpha()]

#         if len(valid_words) < 4:
#             continue

#         eligible_indexes = [i for i, w in enumerate(words) if w.isalpha()]
#         if not eligible_indexes:
#             continue

#         blank_index = random.choice(eligible_indexes)
#         answer = words[blank_index]
#         distractors = [w for w in valid_words if w.lower() != answer.lower()]

#         if len(distractors) < 3:
#             continue

#         words[blank_index] = "_____"
#         question = " ".join(words)
#         options = [answer] + random.sample(distractors, 3)
#         random.shuffle(options)

#         valid_questions.append({
#             "question": question,
#             "options": options,
#             "answer": answer
#         })

#         if len(valid_questions) == num_questions:
#             break

#     if len(valid_questions) < num_questions:
#         st.warning(f"⚠️ Only {len(valid_questions)} MCQs could be generated due to sentence/word limits.")

#     return valid_questions

# # --------------------------- Main App --------------------------- #
# def main():
#     st.set_page_config(page_title="EduLearn - AI Summarizer + MCQ Generator", layout="wide", page_icon="🧠")
#     st.title("🧠 EduLearn")
#     st.subheader("AI Summarizer + MCQ Generator")

#     uploaded_file = st.file_uploader("📂 Upload PDF, DOCX, TXT, or PPTX", type=["pdf", "docx", "txt", "pptx"])

#     if uploaded_file:
#         file_type = uploaded_file.type
#         if "pdf" in file_type:
#             text = read_pdf(uploaded_file)
#         elif "word" in file_type or "docx" in uploaded_file.name:
#             text = read_docx(uploaded_file)
#         elif "text" in file_type:
#             text = read_txt(uploaded_file)
#         elif "pptx" in uploaded_file.name:
#             text = read_pptx(uploaded_file)
#         else:
#             st.error("❌ Unsupported file format")
#             return

#         # --- Summary Section --- #
#         st.header("📚 Summarization")
#         summary_type = st.radio("🔍 Choose summary type:", ["Extractive", "Abstractive"])
#         num_sentences = st.slider("✂ Select number of sentences:", 3, 20, 5)

#         if summary_type == "Extractive":
#             summary = extractive_summary(text, num_sentences)
#         else:
#             method = st.radio("📌 Abstractive Method:", ['lsa', 'lexrank'])
#             st.info("LSA: Best for structured academic text") 
#             st.info("LexRank: Better for long or unstructured content")
#             summary = abstractive_summary(text, num_sentences, method)

#         st.subheader("📝 Summary")
#         st.text_area("Summary Output:", summary, height=200)

#         if st.button("📋 Copy Summary"):
#             pyperclip.copy(summary)
#             st.success("✅ Copied to clipboard!")

#         if st.button("💾 Download Summary"):
#             with open("summary.txt", "w", encoding="utf-8") as f:
#                 f.write(summary)
#             st.download_button(label="📥 Download", data=summary, file_name="summary.txt", mime="text/plain")

#         # --- MCQ Section --- #
#         st.header("❓ MCQ Generator")
#         num_mcqs = st.slider("🔢 Select number of MCQs:", 5, 25, 5)

#         if st.button("🎯 Generate MCQs"):
#             mcqs = generate_mcqs(text, num_questions=num_mcqs)
#             st.session_state["mcqs"] = mcqs
#             st.session_state["show_answers"] = False

#         if "mcqs" in st.session_state:
#             for i, q in enumerate(st.session_state["mcqs"], 1):
#                 st.markdown(f"**Q{i}. {q['question']}**")
#                 st.radio(
#                     f"Choose an option for Q{i}:",
#                     q["options"],
#                     key=f"mcq_{i}_option"
#                 )

#             col1, col2 = st.columns(2)
#             with col1:
#                 if st.button("👁️ Show Answers"):
#                     st.session_state["show_answers"] = True

#             with col2:
#                 if st.button("❌ Hide Answers"):
#                     st.session_state["show_answers"] = False

#             if st.session_state.get("show_answers", False):
#                 st.markdown("### ✅ Correct Answers")
#                 for i, q in enumerate(st.session_state["mcqs"], 1):
#                     st.markdown(f"**Q{i} Answer:** {q['answer']}")

# if __name__ == "__main__":
#     main()


# import streamlit as st
# import PyPDF2
# import docx
# import nltk
# from nltk.tokenize import sent_tokenize
# from sumy.parsers.plaintext import PlaintextParser
# from sumy.nlp.tokenizers import Tokenizer
# from sumy.summarizers.lsa import LsaSummarizer
# from sumy.summarizers.lex_rank import LexRankSummarizer
# from pptx import Presentation
# import random
# import re

# nltk.download('punkt', quiet=True)

# # --------------------------- File Reading --------------------------- #
# def read_pptx(file):
#     prs = Presentation(file)
#     text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
#     return text

# def read_pdf(file):
#     pdf_reader = PyPDF2.PdfReader(file)
#     text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
#     return text

# def read_docx(file):
#     doc = docx.Document(file)
#     text = "\n".join([para.text for para in doc.paragraphs])
#     return text

# def read_txt(file):
#     return file.read().decode("utf-8")

# # --------------------------- Summarization --------------------------- #


# def preprocess_text(text):
#     # Remove numbered lists like 1., 2., 3. etc.
#     text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
#     text = re.sub(r'\n{2,}', ' ', text)  # Replace multiple newlines with space
#     return text.strip()

# def extractive_summary(text, num_sentences):
#     clean_text = preprocess_text(text)
#     sentences = sent_tokenize(clean_text)
#     selected = sentences[:num_sentences]

#     formatted = [f"• {sent.strip()}" for sent in selected if sent.strip()]
#     return "\n".join(formatted)


# def clean_abstractive_output(summary, method):
#     if method == "lsa" or method == "lexrank":
#         # Removing any unwanted characters like '*' and ensuring proper sentence structure
#         summary = summary.replace("*", "")  # Removes asterisks (used in LexRank or LSA for bullet points)
#         summary = re.sub(r'\s+', ' ', summary)  # Removes excessive spaces
#         summary = summary.strip()
#     return summary


# def abstractive_summary(text, num_sentences=5, method="lsa"):
#     parser = PlaintextParser.from_string(text, Tokenizer("english"))
    
#     summarizer = LsaSummarizer() if method == "lsa" else LexRankSummarizer()
    
#     try:
        
#         all_sentences = sent_tokenize(text)
#         summary_sentences = summarizer(parser.document, min(num_sentences * 2, len(all_sentences)))
#         summary_texts = [str(s).strip() for s in summary_sentences]

#         # Preserve original order
#         final_sentences = [s for s in all_sentences if s.strip() in summary_texts][:num_sentences]

#         clean_summary = clean_abstractive_output(" ".join(sentences), method)
#         return clean_summary
#     except Exception as e:
#         st.error(f"Error with {method} summarization: {str(e)}")
#         return ""


# # --------------------------- MCQ Generator --------------------------- #
# def generate_mcqs(text, num_questions=3):
#     sentences = sent_tokenize(text)
#     random.shuffle(sentences)
#     valid_questions = []

#     for sentence in sentences:
#         words = sentence.split()
#         valid_words = [w for w in words if w.isalpha()]

#         if len(valid_words) < 4:
#             continue

#         eligible_indexes = [i for i, w in enumerate(words) if w.isalpha()]
#         if not eligible_indexes:
#             continue

#         blank_index = random.choice(eligible_indexes)
#         answer = words[blank_index]
#         distractors = [w for w in valid_words if w.lower() != answer.lower()]

#         if len(distractors) < 3:
#             continue

#         words[blank_index] = "_____"
#         question = " ".join(words)
#         options = [answer] + random.sample(distractors, 3)
#         random.shuffle(options)

#         valid_questions.append({
#             "question": question,
#             "options": options,
#             "answer": answer
#         })

#         if len(valid_questions) == num_questions:
#             break

#     if len(valid_questions) < num_questions:
#         st.warning(f"⚠️ Only {len(valid_questions)} MCQs could be generated due to sentence/word limits.")

#     return valid_questions

# # --------------------------- Main App --------------------------- #
# def main():
#     st.set_page_config(page_title="EduLearn - AI Summarizer + MCQ Generator", layout="wide", page_icon="🧠")
#     st.title("🧠 EduLearn")
#     st.subheader("AI Summarizer + MCQ Generator")

#     uploaded_file = st.file_uploader("📂 Upload PDF, DOCX, TXT, or PPTX", type=["pdf", "docx", "txt", "pptx"])

#     if uploaded_file:
#         file_type = uploaded_file.type
#         if "pdf" in file_type:
#             text = read_pdf(uploaded_file)
#         elif "word" in file_type or "docx" in uploaded_file.name:
#             text = read_docx(uploaded_file)
#         elif "text" in file_type:
#             text = read_txt(uploaded_file)
#         elif "pptx" in uploaded_file.name:
#             text = read_pptx(uploaded_file)
#         else:
#             st.error("❌ Unsupported file format")
#             return

#         # --- Summary Section --- #
#         st.header("📚 Summarization")
#         summary_type = st.radio("🔍 Choose summary type:", ["Extractive", "Abstractive"])
#         num_sentences = st.slider("✂ Select number of sentences:", 3, 20, 5)

#         if summary_type == "Extractive":
#             summary = extractive_summary(text, num_sentences)
#         else:
#             method = st.radio("📌 Abstractive Method:", ['lsa', 'lexrank'])
#             st.info("LSA: Best for structured academic text") 
#             st.info("LexRank: Better for long or unstructured content")
#             summary = abstractive_summary(text, num_sentences, method)

#         st.subheader("📝 Summary")
#         st.text_area("Summary Output:", summary, height=200)

#         if st.button("💾 Download Summary"):
#             with open("summary.txt", "w", encoding="utf-8") as f:
#                 f.write(summary)
#             st.download_button(label="📥 Download", data=summary, file_name="summary.txt", mime="text/plain")

#         # --- MCQ Section --- #
#         st.header("❓ MCQ Generator")
#         num_mcqs = st.slider("🔢 Select number of MCQs:", 5, 25, 5)

#         if st.button("🎯 Generate MCQs"):
#             mcqs = generate_mcqs(text, num_questions=num_mcqs)
#             st.session_state["mcqs"] = mcqs
#             st.session_state["show_answers"] = False

#         if "mcqs" in st.session_state:
#             for i, q in enumerate(st.session_state["mcqs"], 1):
#                 st.markdown(f"**Q{i}. {q['question']}**")
#                 st.radio(f"Choose an option for Q{i}:", q["options"], key=f"mcq_{i}_option")

#             col1, col2 = st.columns(2)
#             with col1:
#                 if st.button("👁️ Show Answers"):
#                     st.session_state["show_answers"] = True

#             with col2:
#                 if st.button("❌ Hide Answers"):
#                     st.session_state["show_answers"] = False

#             if st.session_state.get("show_answers", False):
#                 st.markdown("### ✅ Correct Answers")
#                 for i, q in enumerate(st.session_state["mcqs"], 1):
#                     st.markdown(f"**Q{i} Answer:** {q['answer']}")

# if __name__ == "__main__":
#     main()




import streamlit as st
import PyPDF2
import docx
import nltk
from nltk.tokenize import sent_tokenize
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from pptx import Presentation
import random
import re

nltk.download('punkt', quiet=True)

# --------------------------- File Reading --------------------------- #
def read_pptx(file):
    prs = Presentation(file)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

def read_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
    return text

def read_docx(file):
    doc = docx.Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def read_txt(file):
    return file.read().decode("utf-8")

# --------------------------- Summarization --------------------------- #
def preprocess_text(text):
    text = re.sub(r'^\s*\d+\.\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\n{2,}', ' ', text)
    return text.strip()

def extractive_summary(text, num_sentences):
    clean_text = preprocess_text(text)
    sentences = sent_tokenize(clean_text)
    selected = sentences[:num_sentences]
    formatted = [f"• {sent.strip()}" for sent in selected if sent.strip()]
    return "\n".join(formatted)

def clean_abstractive_output(summary, method):
    if method in ["lsa", "lexrank"]:
        summary = summary.replace("*", "")
        summary = re.sub(r'\s+', ' ', summary)
        summary = summary.strip()
    return summary


def abstractive_summary(text, num_sentences=5, method="lsa"):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer() if method == "lsa" else LexRankSummarizer()

    try:
        all_sentences = sent_tokenize(text)
        summary_sentences = summarizer(parser.document, min(num_sentences * 2, len(all_sentences)))
        summary_texts = [str(s).strip() for s in summary_sentences]

        # Filter from original sentences for better coherence
        final_sentences = [s for s in all_sentences if s.strip() in summary_texts][:num_sentences]

        # Clean & return
        cleaned = [s for s in final_sentences if s and not re.fullmatch(r"\d+\.?", s.strip())]
        return " ".join(cleaned) if method == "lexrank" else "\n".join(cleaned)
    except Exception as e:
        st.error(f"Error in {method.upper()} summary: {e}")
        return ""

# --------------------------- Formatted Summary Display --------------------------- #

def formatted_summary(summary_text):
    sentences = sent_tokenize(summary_text)
    numbered = []
    for i, s in enumerate(sentences):
        s = s.strip()
        if not s or re.fullmatch(r"\d+\.?", s):  # skip empty or numeric-only lines
            continue
        numbered.append(f"{len(numbered)+1}. {s}")
    return "\n".join(numbered)

# --------------------------- MCQ Generator --------------------------- #
def generate_mcqs(text, num_questions=3):
    sentences = sent_tokenize(text)
    random.shuffle(sentences)
    valid_questions = []

    for sentence in sentences:
        words = sentence.split()
        valid_words = [w for w in words if w.isalpha()]

        if len(valid_words) < 4:
            continue

        eligible_indexes = [i for i, w in enumerate(words) if w.isalpha()]
        if not eligible_indexes:
            continue

        blank_index = random.choice(eligible_indexes)
        answer = words[blank_index]
        distractors = [w for w in valid_words if w.lower() != answer.lower()]

        if len(distractors) < 3:
            continue

        words[blank_index] = "_____"
        question = " ".join(words)
        options = [answer] + random.sample(distractors, 3)
        random.shuffle(options)

        valid_questions.append({
            "question": question,
            "options": options,
            "answer": answer
        })

        if len(valid_questions) == num_questions:
            break

    if len(valid_questions) < num_questions:
        st.warning(f"⚠️ Only {len(valid_questions)} MCQs could be generated due to sentence/word limits.")

    return valid_questions


# --------------------------- Main App --------------------------- #
def main():
    st.set_page_config(page_title="EduLearn - AI Summarizer + MCQ Generator", layout="wide", page_icon="🧠")
    st.title("🧠 EduLearn")
    st.subheader("AI Summarizer + MCQ Generator")

    uploaded_file = st.file_uploader("📂 Upload PDF, DOCX, TXT, or PPTX", type=["pdf", "docx", "txt", "pptx"])

    if uploaded_file:
        file_type = uploaded_file.type
        if "pdf" in file_type:
            text = read_pdf(uploaded_file)
        elif "word" in file_type or "docx" in uploaded_file.name:
            text = read_docx(uploaded_file)
        elif "text" in file_type:
            text = read_txt(uploaded_file)
        elif "pptx" in uploaded_file.name:
            text = read_pptx(uploaded_file)
        else:
            st.error("❌ Unsupported file format")
            return

        # --- Summary Section --- #
        st.header("📚 Summarization")
        summary_type = st.radio("🔍 Choose summary type:", ["Extractive", "Abstractive"])
        num_sentences = st.slider("✂ Select number of sentences:", 3, 20, 5)

        if summary_type == "Extractive":
            summary = extractive_summary(text, num_sentences)
        else:
            method = st.radio("📌 Abstractive Method:", ['lsa', 'lexrank'])
            st.info("LSA: Best for structured academic text") 
            st.info("LexRank: Better for long or unstructured content")
            summary = abstractive_summary(text, num_sentences, method)

        st.subheader("📝 Summary")
        show_formatted = st.checkbox("🧾 Show formatted (numbered) summary", value=False)

        # Control formatting based on method and toggle
        if summary_type == "Extractive":
            display_summary = formatted_summary(summary) if show_formatted else summary
        else:  # Abstractive
            if show_formatted and method != "lexrank":
                display_summary = formatted_summary(summary)
            else:
                display_summary = summary  # LexRank stays clean paragraph

        st.text_area("Summary Output:", display_summary, height=250)

        if st.button("💾 Download Summary"):
            st.download_button(label="📥 Download", data=display_summary, file_name="summary.txt", mime="text/plain")

        # --- MCQ Section --- #
        st.header("❓ MCQ Generator")
        num_mcqs = st.slider("🔢 Select number of MCQs:", 5, 25, 5)

        if st.button("🎯 Generate MCQs"):
            mcqs = generate_mcqs(text, num_questions=num_mcqs)
            st.session_state["mcqs"] = mcqs
            st.session_state["show_answers"] = False

        if "mcqs" in st.session_state:
            for i, q in enumerate(st.session_state["mcqs"], 1):
                st.markdown(f"**Q{i}. {q['question']}**")
                st.radio(f"Choose an option for Q{i}:", q["options"], key=f"mcq_{i}_option")

            col1, col2 = st.columns(2)
            with col1:
                if st.button("👁️ Show Answers"):
                    st.session_state["show_answers"] = True

            with col2:
                if st.button("❌ Hide Answers"):
                    st.session_state["show_answers"] = False

            if st.session_state.get("show_answers", False):
                st.markdown("### ✅ Correct Answers")
                for i, q in enumerate(st.session_state["mcqs"], 1):
                    st.markdown(f"**Q{i} Answer:** {q['answer']}")

if __name__ == "__main__":
    main()
